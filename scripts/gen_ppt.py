from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
AZUL = RGBColor(0x17, 0x4B, 0x87)
AZUL_CLARO = RGBColor(0x5B, 0x9B, 0xD5)
GRIS = RGBColor(0x44, 0x44, 0x44)


def nueva_diapositiva_titulo(prs, titulo, subtitulo=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ancho, alto = prs.slide_width, prs.slide_height

    barra = slide.shapes.add_shape(1, 0, 0, ancho, Inches(1.2))
    barra.fill.solid()
    barra.fill.fore_color.rgb = AZUL
    barra.line.fill.background()

    caja = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), ancho - Inches(1), Inches(0.8))
    tf = caja.text_frame
    p = tf.paragraphs[0]
    p.text = titulo
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = BLANCO

    if subtitulo:
        caja2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), ancho - Inches(1), Inches(0.6))
        tf2 = caja2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitulo
        p2.font.size = Pt(16)
        p2.font.color.rgb = GRIS
    return slide


def nueva_diapositiva_codigo(prs, titulo):
    slide = prs.slides.add_slide(prs.slide_layouts[4])
    slide.shapes.title.text = titulo
    return slide


def añadir_bullets(slide, items, monto=Inches(0.3)):
    caja = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.3), Inches(5.5))
    tf = caja.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            texto, nivel = item
        else:
            texto, nivel = item, 0
        p.text = texto
        p.level = nivel
        p.font.size = Pt(18 if nivel == 0 else 15)
        p.font.color.rgb = GRIS
        p.space_after = Pt(6)
    return caja


def añadir_tabla(slide, datos, filas_columnas=None, col_widths=None):
    filas = len(datos)
    columnas = len(datos[0])
    tabla = slide.shapes.add_table(filas, columnas, Inches(0.6), Inches(1.5), Inches(12.3), Inches(5.5)).table
    if col_widths:
        for i, w in enumerate(col_widths):
            tabla.columns[i].width = w
    for r, fila in enumerate(datos):
        for c, valor in enumerate(fila):
            celda = tabla.cell(r, c)
            celda.text = str(valor)
            for p in celda.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = BLANCO if r == 0 else GRIS
            celda.text_frame.paragraphs[0].font.bold = r == 0
            if r == 0:
                celda.fill.solid()
                celda.fill.fore_color.rgb = AZUL
    return tabla


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Portada
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "ADAII UT2 TFU"
    slide.placeholders[1].text = (
        "Tácticas de Arquitectura\n"
        "API de compras de juegos  •  FastAPI + MySQL + Docker\n"
        "Arquitectura de Aplicaciones II"
    )

    # 2. Descripción del proyecto
    s = nueva_diapositiva_titulo(prs, "Descripción del proyecto")
    añadir_bullets(s, [
        "Trabajo Final de Unidad (TFU) de la asignatura Arquitectura de Aplicaciones II (UT2).",
        "API de gestión de compras de juegos construida con FastAPI, MySQL y Docker.",
        "Demuestra tácticas de arquitectura aplicadas a un caso real:",
        ("Autenticación JWT con contraseñas hasheadas (bcrypt)", 1),
        ("Patrón Strategy para políticas de costo de compra", 1),
        ("Auditoría de operaciones mediante decoradores (AOP)", 1),
        ("Monitoreo de disponibilidad vía healthcheck", 1),
        ("Despliegue y rollback versionado por contenedores", 1),
    ], Inches(0.3))

    # 3. Stack tecnológico
    s = nueva_diapositiva_titulo(prs, "Stack tecnológico")
    añadir_tabla(s, [
        ["Componente", "Tecnología"],
        ["Lenguaje", "Python 3.12"],
        ["Framework API", "FastAPI + Uvicorn"],
        ["ORM", "SQLAlchemy + PyMySQL"],
        ["Validación", "Pydantic / pydantic-settings"],
        ["Base de datos", "MySQL 8.4"],
        ["Autenticación", "JWT (HS256) + bcrypt"],
        ["Migraciones", "Alembic"],
        ["Contenedores", "Docker / Docker Compose"],
    ], col_widths=[Inches(4), Inches(8)])

    # 4. Tácticas: resumen
    s = nueva_diapositiva_titulo(prs, "Tácticas de arquitectura", "Cinco tácticas implementadas y demostradas")
    añadir_tabla(s, [
        ["Táctica", "Implementación"],
        ["Seguridad — autenticación", "JWT firmado y contraseñas con bcrypt (security/)"],
        ["Abstracción de servicios — Strategy", "Costo de compra según política: normal / invierno"],
        ["Auditoría — AOP", "Decorador @auditar que registra cada operación"],
        ["Disponibilidad — monitorización", "Endpoints /health y /version para detectar fallos"],
        ["Recuperación — rollback", "Imágenes versionadas y reversión ante healthcheck en fallo"],
    ], col_widths=[Inches(4.6), Inches(7.4)])

    # 5. Seguridad / autenticación
    s = nueva_diapositiva_titulo(prs, "Táctica: Seguridad — Autenticación")
    añadir_bullets(s, [
        "Login con JWT firmado por secreto (HS256), expiración configurable.",
        "Contraseñas almacenadas con hash bcrypt (nunca en texto plano).",
        "Dependencia obtener_usuario_actual aplicada por router:",
        ("/personas, /juegos, /compras y /auditoria exigen token válido", 1),
        ("Esquema: Authorization: Bearer <token>", 1),
        "Time to live del token configurable vía JWT_EXPIRATION_MINUTES.",
    ], Inches(0.3))

    # 6. Strategy
    s = nueva_diapositiva_titulo(prs, "Táctica: Strategy — Políticas de costo")
    añadir_bullets(s, [
        "El cálculo del costo de una compra se delega en una estrategia intercambiable.",
        "app/strategies/costo_compra.py define la interfaz EstrategiaCostoCompra:",
        ("normal → CostoNormal: devuelve el costo base", 1),
        ("invierno → DescuentoInvierno: 50% del costo base", 1),
        "Nueva política = nueva clase, sin tocar la lógica de compras.",
        "Política no soportada → ValueError controlado.",
    ], Inches(0.3))

    # 7. Auditoría / AOP
    s = nueva_diapositiva_titulo(prs, "Táctica: Auditoría — AOP")
    añadir_bullets(s, [
        "Decorador @auditar(\"OPERACION\") registra cada operación en un log.",
        "Soporta funciones síncronas y asíncronas (inspect.iscoroutinefunction).",
        "Escribir en archivo con timestamp ISO-8601 UTC, protegido por lock.",
        "Registra resultado OK o ERROR:tipo en cada operación.",
        "Visible en /auditoria (últimas 100 líneas) y ruta AUDIT_LOG_PATH.",
        "Ejemplo: GENERAR_COMPRA, ACTUALIZAR_COMPRA, LOGIN, INSERTAR_PERSONA.",
    ], Inches(0.3))

    # 8. Disponibilidad
    s = nueva_diapositiva_titulo(prs, "Táctica: Disponibilidad — Healthcheck")
    añadir_bullets(s, [
        "Endpoint GET /health devuelve {status: ok}.",
        "FORCE_UNHEALTHY=true simula una versión defectuosa (HTTP 503).",
        "Docker verifica la salud del contenedor cada 5s (mysqladmin ping en BD).",
        "GET /version expone versión desplegada y estado de la instancia.",
        "Permite detectar y excluir instancias degradadas del servicio.",
    ], Inches(0.3))

    # 9. Rollback
    s = nueva_diapositiva_titulo(prs, "Táctica: Recuperación — Rollback versionado")
    añadir_bullets(s, [
        "Cada versión se construye como imagen versionada: api:<version>.",
        "build_version.sh <version> [force_unhealthy] genera la imagen.",
        "deploy.sh guarda la versión anterior (.deploy/current_version) y registra logs.",
        "demo_rollback.sh automatiza la demo completa:",
        ("Despliega 1.0 (estable) → revisa salud", 1),
        ("Despliega 2.0-broken (FORCE_UNHEALTHY=true) → healthcheck falla", 1),
        ("Revierte automáticamente a la última versión estable", 1),
    ], Inches(0.3))

    # 10. Estructura
    s = nueva_diapositiva_titulo(prs, "Estructura del proyecto")
    añadir_bullets(s, [
        ("app/", 0),
        ("main.py — app FastAPI, routers, /health, /version", 1),
        ("models.py / schemas.py — ORM y Pydantic", 1),
        ("aspects/audit.py — decorador @auditar (AOP)", 1),
        ("routers/ — auth, personas, juegos, compras, auditoria", 1),
        ("services/compras.py — lógica de negocio", 1),
        ("strategies/costo_compra.py — patrón Strategy", 1),
        ("bd/ — schema.sql y seed.sql (init automático)", 0),
        ("security/ — jwt.py, password.py, dependencies.py", 0),
        ("scripts/ — .sh y .ps1 de despliegue y demo", 0),
        ("postman/ — colección y entorno de pruebas", 0),
    ], Inches(0.3))

    # 11. Puesta en marcha
    s = nueva_diapositiva_titulo(prs, "Puesta en marcha")
    añadir_bullets(s, [
        "Producción:",
        ("cp .env.example .env", 1),
        ("./scripts/start.sh  (construye api:1.0 y levanta todo)", 1),
        ("API en http://localhost:8000  •  Swagger en /docs", 1),
        ("Usuario demo: demo@adaii.local / Demo123!", 1),
        "Desarrollo (hot reload):",
        ("docker compose -f docker-compose.dev.yaml up --build", 1),
        ("Monta el código local y usa uvicorn --reload", 1),
        ("Nota: MySQL expuesto en el host por el puerto 3307", 0),
    ], Inches(0.3))

    # 12. Scripts
    s = nueva_diapositiva_titulo(prs, "Scripts disponibles", "Versiones .sh (shell) y .ps1 (PowerShell)")
    añadir_tabla(s, [
        ["Script", "Descripción"],
        ["start.sh", "Crea .env, construye api:1.0 y levanta la API"],
        ["build_version.sh <v> [fu]", "Construye imagen versionada; force_unhealthy=true simula fallo"],
        ["deploy.sh <v>", "Despliega una imagen construida y registra el log"],
        ["demo_rollback.sh", "Demo: 1.0 estable → 2.0-broken → rollback automático"],
        ["demo_tacticas.sh", "Demo de tácticas: login JWT y estrategias de costo"],
    ], col_widths=[Inches(4.6), Inches(7.4)], )

    # 13. Endpoints
    s = nueva_diapositiva_titulo(prs, "Endpoints principales")
    añadir_tabla(s, [
        ["Método", "Ruta", "Descripción"],
        ["GET", "/", "Mensaje y versión de la API"],
        ["GET", "/health", "Estado de salud (healthcheck de Docker)"],
        ["GET", "/version", "Versión y estado defectuoso"],
        ["POST", "/auth/login", "Autenticación, devuelve JWT"],
        ["POST", "/personas", "Crear persona (protegido)"],
        ["…", "/juegos, /compras, /auditoria", "CRUD y consultas protegidas con JWT"],
    ], col_widths=[Inches(1.4), Inches(4.4), Inches(6.5)])

    # 14. Variables de entorno
    s = nueva_diapositiva_titulo(prs, "Variables de entorno", "Ver .env.example")
    añadir_tabla(s, [
        ["Variable", "Descripción"],
        ["JWT_SECRET", "Clave secreta para firmar tokens"],
        ["JWT_ALGORITHM / JWT_EXPIRATION_MINUTES", "Algoritmo y validez del token"],
        ["MYSQL_ROOT_PASSWORD", "Password root de MySQL"],
        ["MYSQL_DATABASE / MYSQL_USER / MYSQL_PASSWORD", "BD y credenciales de la app"],
        ["MYSQL_HOST / MYSQL_PORT", "Seteados por Docker Compose"],
        ["AUDIT_LOG_PATH", "Ruta del log de auditoría"],
    ], col_widths=[Inches(6), Inches(6.3)])

    # 15. Cierre
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    caja = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
    tf = caja.text_frame
    p = tf.paragraphs[0]
    p.text = "Gracias"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = AZUL
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "Tácticas de arquitectura demostradas en una API real, lista para ejecutar y auditar."
    p2.font.size = Pt(18)
    p2.font.color.rgb = GRIS
    p2.alignment = PP_ALIGN.CENTER

    prs.save("ADAII-UT2_Presentacion.pptx")
    print("OK: ADAII-UT2_Presentacion.pptx")


if __name__ == "__main__":
    main()